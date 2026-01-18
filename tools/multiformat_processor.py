"""
MULTIFORMAT PROCESSOR - Xử lý đa định dạng tài liệu
"""
import os
import sys
import json
import hashlib
from pathlib import Path
from datetime import datetime
from typing import Dict, Any, List, Optional
import logging

# Thêm đường dẫn
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

class MultiFormatProcessor:
    """Xử lý đa định dạng tài liệu"""
    
    def __init__(self):
        self.logger = logging.getLogger(__name__)
        self.supported_formats = {
            # Text formats
            '.txt': self._process_text,
            '.md': self._process_markdown,
            '.json': self._process_json,
            '.yaml': self._process_yaml,
            '.yml': self._process_yaml,
            '.csv': self._process_csv,
            '.xml': self._process_xml,
            '.html': self._process_html,
            '.htm': self._process_html,
            
            # Office formats
            '.xlsx': self._process_excel,
            '.xls': self._process_excel,
            '.pptx': self._process_powerpoint,
            '.ppt': self._process_powerpoint,
            '.docx': self._process_word,
            '.doc': self._process_word_legacy,
            
            # PDF
            '.pdf': self._process_pdf,
            
            # Images
            '.jpg': self._process_image,
            '.jpeg': self._process_image,
            '.png': self._process_image,
            '.bmp': self._process_image,
            '.gif': self._process_image,
            '.tiff': self._process_image,
            
            # Code
            '.py': self._process_python,
            '.js': self._process_code,
            '.java': self._process_code,
            '.cpp': self._process_code,
            '.c': self._process_code,
            '.html': self._process_code,
            '.css': self._process_code,
        }
    
    def process_file(self, file_path: str) -> Dict[str, Any]:
        """Xử lý file đa định dạng"""
        try:
            file_path = Path(file_path)
            
            if not file_path.exists():
                return {"error": f"File không tồn tại: {file_path}"}
            
            # Kiểm tra định dạng
            extension = file_path.suffix.lower()
            if extension not in self.supported_formats:
                return {"error": f"Định dạng không hỗ trợ: {extension}"}
            
            # Xử lý file
            processor = self.supported_formats[extension]
            result = processor(file_path)
            
            # Thêm metadata
            result.update({
                "filename": file_path.name,
                "filepath": str(file_path),
                "extension": extension,
                "file_size": file_path.stat().st_size,
                "modified_time": datetime.fromtimestamp(file_path.stat().st_mtime).isoformat(),
                "processing_time": datetime.now().isoformat(),
                "content_hash": hashlib.md5(str(result).encode()).hexdigest()[:8]
            })
            
            return result
            
        except Exception as e:
            self.logger.error(f"Lỗi xử lý file {file_path}: {e}")
            return {"error": str(e), "filename": file_path.name}
    
    def _process_text(self, file_path: Path) -> Dict[str, Any]:
        """Xử lý file text"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        return {
            "content": content,
            "type": "text",
            "word_count": len(content.split()),
            "line_count": len(content.splitlines()),
            "encoding": "utf-8"
        }
    
    def _process_markdown(self, file_path: Path) -> Dict[str, Any]:
        """Xử lý markdown"""
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Phân tích markdown cơ bản
        lines = content.splitlines()
        headers = [line for line in lines if line.startswith('#')]
        code_blocks = content.count('```')
        
        return {
            "content": content,
            "type": "markdown",
            "headers": len(headers),
            "code_blocks": code_blocks,
            "word_count": len(content.split())
        }
    
    def _process_json(self, file_path: Path) -> Dict[str, Any]:
        """Xử lý JSON"""
        import json as json_module
        
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json_module.load(f)
        
        return {
            "content": json_module.dumps(data, indent=2),
            "type": "json",
            "data_structure": self._analyze_json_structure(data),
            "is_valid": True
        }
    
    def _analyze_json_structure(self, data: Any) -> Dict[str, Any]:
        """Phân tích cấu trúc JSON"""
        if isinstance(data, dict):
            return {
                "type": "object",
                "keys": list(data.keys()),
                "key_count": len(data)
            }
        elif isinstance(data, list):
            return {
                "type": "array",
                "length": len(data),
                "first_item_type": type(data[0]).__name__ if data else None
            }
        else:
            return {"type": type(data).__name__}
    
    def _process_excel(self, file_path: Path) -> Dict[str, Any]:
        """Xử lý Excel files"""
        try:
            import pandas as pd
            
            # Đọc tất cả sheets
            excel_file = pd.ExcelFile(file_path)
            sheets = {}
            
            for sheet_name in excel_file.sheet_names:
                try:
                    df = pd.read_excel(excel_file, sheet_name=sheet_name)
                    
                    # Convert to text representation
                    sheet_content = {
                        "headers": df.columns.tolist(),
                        "row_count": len(df),
                        "column_count": len(df.columns),
                        "sample_data": df.head(5).to_dict(orient='records'),
                        "dtypes": {col: str(dtype) for col, dtype in df.dtypes.items()}
                    }
                    
                    sheets[sheet_name] = sheet_content
                    
                except Exception as e:
                    sheets[sheet_name] = {"error": str(e)}
            
            # Tạo text summary
            text_content = f"EXCEL FILE: {file_path.name}\n\n"
            for sheet_name, data in sheets.items():
                if "error" not in data:
                    text_content += f"SHEET: {sheet_name}\n"
                    text_content += f"Rows: {data['row_count']}, Columns: {data['column_count']}\n"
                    text_content += f"Headers: {', '.join(data['headers'][:5])}\n\n"
            
            return {
                "content": text_content,
                "type": "excel",
                "sheets": sheets,
                "sheet_count": len(sheets),
                "processed_with": "pandas"
            }
            
        except ImportError:
            return {"error": "Thư viện pandas chưa được cài đặt"}
        except Exception as e:
            return {"error": f"Lỗi xử lý Excel: {e}"}
    
    def _process_powerpoint(self, file_path: Path) -> Dict[str, Any]:
        """Xử lý PowerPoint"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            slides_content = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                
                # Lấy text từ shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)
                
                slides_content.append({
                    "slide_number": i + 1,
                    "text_content": slide_text,
                    "text_count": len(slide_text)
                })
            
            # Tạo text summary
            text_content = f"POWERPOINT: {file_path.name}\n\n"
            for slide in slides_content:
                text_content += f"Slide {slide['slide_number']}:\n"
                for text in slide['text_content'][:3]:  # Lấy 3 dòng đầu
                    text_content += f"  • {text}\n"
                text_content += "\n"
            
            return {
                "content": text_content,
                "type": "powerpoint",
                "slides": slides_content,
                "slide_count": len(slides_content),
                "total_text_slides": sum(s['text_count'] > 0 for s in slides_content)
            }
            
        except ImportError:
            return {"error": "Thư viện python-pptx chưa được cài đặt"}
        except Exception as e:
            return {"error": f"Lỗi xử lý PowerPoint: {e}"}
    
    def _process_word(self, file_path: Path) -> Dict[str, Any]:
        """Xử lý Word (.docx)"""
        try:
            from docx import Document
            
            doc = Document(file_path)
            paragraphs = []
            tables_data = []
            
            # Lấy text từ paragraphs
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs.append(text)
            
            # Lấy text từ tables
            for table in doc.tables:
                table_content = []
                for row in table.rows:
                    row_content = [cell.text.strip() for cell in row.cells]
                    table_content.append(row_content)
                tables_data.append(table_content)
            
            # Tạo text summary
            text_content = f"WORD DOCUMENT: {file_path.name}\n\n"
            text_content += "\n".join(paragraphs[:20])  # Lấy 20 đoạn đầu
            
            return {
                "content": text_content,
                "type": "word",
                "full_text": "\n".join(paragraphs),
                "paragraph_count": len(paragraphs),
                "table_count": len(tables_data),
                "word_count": sum(len(p.split()) for p in paragraphs)
            }
            
        except ImportError:
            return {"error": "Thư viện python-docx chưa được cài đặt"}
        except Exception as e:
            return {"error": f"Lỗi xử lý Word: {e}"}
    
    def _process_word_legacy(self, file_path: Path) -> Dict[str, Any]:
        """Xử lý Word cũ (.doc) - cần antiword hoặc convert"""
        return {
            "content": f"Legacy Word file: {file_path.name}\nCần chuyển đổi sang .docx để đọc nội dung.",
            "type": "word_legacy",
            "needs_conversion": True
        }
    
    def _process_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Xử lý PDF files"""
        try:
            # Thử PyPDF2 trước
            import PyPDF2
            
            text_content = ""
            page_count = 0
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                page_count = len(pdf_reader.pages)
                
                for page_num in range(min(10, page_count)):  # Đọc tối đa 10 trang
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    if text:
                        text_content += f"--- Trang {page_num + 1} ---\n{text}\n\n"
            
            # Nếu không có text, thử pdfminer
            if not text_content.strip():
                try:
                    from pdfminer.high_level import extract_text
                    text_content = extract_text(file_path)
                except:
                    text_content = "Không thể extract text từ PDF"
            
            return {
                "content": text_content[:10000],  # Giới hạn nội dung
                "type": "pdf",
                "page_count": page_count,
                "text_extracted": len(text_content) > 100,
                "text_length": len(text_content)
            }
            
        except ImportError:
            return {"error": "Thư viện PyPDF2/pdfminer chưa được cài đặt"}
        except Exception as e:
            return {"error": f"Lỗi xử lý PDF: {e}"}
    
    def _process_image(self, file_path: Path) -> Dict[str, Any]:
        """Xử lý hình ảnh với OCR"""
        try:
            from PIL import Image
            import pytesseract
            
            # Mở hình ảnh
            image = Image.open(file_path)
            
            # Get image info
            image_info = {
                "format": image.format,
                "size": image.size,
                "mode": image.mode,
                "width": image.width,
                "height": image.height
            }
            
            # Thử OCR
            try:
                text = pytesseract.image_to_string(image, lang='vie+eng')
                ocr_success = len(text.strip()) > 10
            except:
                text = "Không thể thực hiện OCR. Cần cài đặt Tesseract."
                ocr_success = False
            
            return {
                "content": text,
                "type": "image",
                "image_info": image_info,
                "ocr_applied": True,
                "ocr_success": ocr_success,
                "text_from_image": text[:5000] if ocr_success else "N/A"
            }
            
        except ImportError as e:
            return {"error": f"Thư viện chưa được cài đặt: {e}"}
        except Exception as e:
            return {"error": f"Lỗi xử lý hình ảnh: {e}"}
    
    def _process_python(self, file_path: Path) -> Dict[str, Any]:
        """Xử lý Python code"""
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Phân tích code cơ bản
        lines = content.splitlines()
        imports = [line for line in lines if line.strip().startswith('import') or line.strip().startswith('from')]
        functions = [line for line in lines if line.strip().startswith('def ')]
        classes = [line for line in lines if line.strip().startswith('class ')]
        
        return {
            "content": content,
            "type": "python_code",
            "line_count": len(lines),
            "import_count": len(imports),
            "function_count": len(functions),
            "class_count": len(classes),
            "sample_functions": [f.split('def ')[1].split('(')[0] for f in functions[:3]]
        }
    
    def _process_csv(self, file_path: Path) -> Dict[str, Any]:
        """Xử lý CSV"""
        try:
            import pandas as pd
            
            df = pd.read_csv(file_path)
            
            return {
                "content": df.head(20).to_string(),
                "type": "csv",
                "row_count": len(df),
                "column_count": len(df.columns),
                "headers": df.columns.tolist(),
                "sample_data": df.head(5).to_dict(orient='records')
            }
            
        except Exception as e:
            return {"error": f"Lỗi xử lý CSV: {e}"}
    
    def _process_yaml(self, file_path: Path) -> Dict[str, Any]:
        """Xử lý YAML"""
        try:
            import yaml
            
            with open(file_path, 'r', encoding='utf-8') as f:
                data = yaml.safe_load(f)
            
            return {
                "content": yaml.dump(data, default_flow_style=False),
                "type": "yaml",
                "is_valid": True,
                "data_type": type(data).__name__
            }
            
        except Exception as e:
            return {"error": f"Lỗi xử lý YAML: {e}"}
    
    def _process_xml(self, file_path: Path) -> Dict[str, Any]:
        """Xử lý XML"""
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        return {
            "content": content[:5000],  # Giới hạn
            "type": "xml",
            "has_xml_structure": "<?xml" in content or "<root>" in content
        }
    
    def _process_html(self, file_path: Path) -> Dict[str, Any]:
        """Xử lý HTML"""
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Extract text từ HTML đơn giản
        import re
        text_content = re.sub('<[^<]+?>', ' ', content)
        text_content = re.sub('\s+', ' ', text_content).strip()
        
        return {
            "content": text_content[:5000],
            "type": "html",
            "original_html": content[:2000],
            "tag_count": content.count('<'),
            "text_extracted": text_content[:500]
        }
    
    def _process_code(self, file_path: Path) -> Dict[str, Any]:
        """Xử lý code files khác"""
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        return {
            "content": content,
            "type": "code",
            "line_count": len(content.splitlines()),
            "file_type": file_path.suffix
        }
    
    def batch_process(self, folder_path: str, extensions: List[str] = None) -> Dict[str, Any]:
        """Xử lý hàng loạt"""
        folder = Path(folder_path)
        results = {
            "total_files": 0,
            "processed": 0,
            "failed": 0,
            "results": [],
            "summary_by_type": {}
        }
        
        if not folder.exists():
            return {"error": f"Thư mục không tồn tại: {folder_path}"}
        
        # Lấy tất cả file
        files = []
        for ext in (extensions or self.supported_formats.keys()):
            files.extend(folder.rglob(f"*{ext}"))
        
        for file_path in files:
            if file_path.is_file():
                results["total_files"] += 1
                
                try:
                    file_result = self.process_file(str(file_path))
                    
                    if "error" not in file_result:
                        results["processed"] += 1
                        
                        # Thống kê theo type
                        file_type = file_result.get("type", "unknown")
                        if file_type not in results["summary_by_type"]:
                            results["summary_by_type"][file_type] = 0
                        results["summary_by_type"][file_type] += 1
                        
                        results["results"].append({
                            "file": file_path.name,
                            "type": file_type,
                            "success": True,
                            "size": file_result.get("file_size", 0),
                            "preview": str(file_result.get("content", ""))[:100]
                        })
                    else:
                        results["failed"] += 1
                        results["results"].append({
                            "file": file_path.name,
                            "error": file_result["error"],
                            "success": False
                        })
                        
                except Exception as e:
                    results["failed"] += 1
                    results["results"].append({
                        "file": file_path.name,
                        "error": str(e),
                        "success": False
                    })
        
        return results